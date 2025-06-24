
import streamlit as st
import lyricsgenius
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import io

# Streamlit UI
st.title("🎤 Lyrics to PowerPoint")
st.write("Zadej interpreta a název písně. Aplikace stáhne text a vytvoří PowerPoint prezentaci.")

# User inputs
artist = st.text_input("Interpret")
title = st.text_input("Název písně")

# Function to get lyrics
def get_lyrics(artist, title, token):
    genius = lyricsgenius.Genius(token, timeout=15, retries=3)
    song = genius.search_song(title, artist)
    return song.lyrics if song and song.lyrics else None

# Function to divide lyrics into chunks
def divide_text(lyrics):
    lines = [l for l in lyrics.splitlines() if l.strip()]
    return ["\n".join(lines[i:i+4]) for i in range(0, len(lines), 4)]

# Function to create PowerPoint
def create_ppt(chunks, song_title, artist):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]

    # Title slide
    slide = prs.slides.add_slide(layout)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1))
    title_tf = title_box.text_frame
    title_tf.text = song_title
    title_tf.paragraphs[0].font.size = Pt(54)
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(1))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.text = artist
    subtitle_tf.paragraphs[0].font.size = Pt(32)
    subtitle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    for chunk in chunks:
        slide = prs.slides.add_slide(layout)
        box = slide.shapes.add_textbox(Inches(1.5), Inches(1), Inches(10), Inches(5.5))
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        for i, line in enumerate(chunk.splitlines()):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(32)
            p.alignment = PP_ALIGN.CENTER

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# Button to generate presentation
if st.button("🎵 Vytvořit prezentaci"):
    if not artist or not title:
        st.warning("Zadej interpreta i název písně.")
    else:
        GENIUS_TOKEN = st.secrets["GENIUS_TOKEN"]
        with st.spinner("🔍 Hledám text písně..."):
            lyrics = get_lyrics(artist, title, GENIUS_TOKEN)
        if lyrics:
            chunks = divide_text(lyrics)
            ppt_file = create_ppt(chunks, title, artist)
            filename = f"lyrics_{title}_by_{artist}.pptx".replace(" ", "_")
            st.success("✅ Prezentace vytvořena!")
            st.download_button("📥 Stáhnout prezentaci", ppt_file, file_name=filename)
        else:
            st.error("❌ Text písně nebyl nalezen.")
